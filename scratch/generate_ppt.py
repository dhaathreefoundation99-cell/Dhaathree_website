import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ----------------------------------------------------
    # Color Palette & Styling Settings
    # ----------------------------------------------------
    NAVY = RGBColor(12, 35, 64)       # Primary (Adani Navy)
    ORANGE = RGBColor(242, 101, 34)    # Accent (Safety Orange)
    DARK_GRAY = RGBColor(60, 60, 60)   # Body Text
    LIGHT_GRAY = RGBColor(240, 242, 245) # Box Backgrounds
    WHITE = RGBColor(255, 255, 255)
    
    # Helper to add standard header and background shapes
    def apply_slide_template(slide, title_text, subtitle_text=None):
        # Header background
        header_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = NAVY
        header_shape.line.fill.background()
        
        # Title text box
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(12.333), Inches(0.9))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Segoe UI"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.name = "Segoe UI"
            p2.font.size = Pt(14)
            p2.font.color.rgb = ORANGE
            p2.font.italic = True
            
        # Top banner separator line
        line_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.1), Inches(13.333), Inches(0.08)
        )
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = ORANGE
        line_shape.line.fill.background()

    def add_bullet_points(tf, points, font_size=13, space_after=12):
        for i, pt in enumerate(points):
            p = tf.add_paragraph() if i > 0 or tf.paragraphs[0].text else tf.paragraphs[0]
            p.text = ""
            p.font.name = "Segoe UI"
            p.font.size = Pt(font_size)
            p.space_after = Pt(space_after)
            
            # Formats first words as bold if they contain a colon
            if ":" in pt:
                lead, detail = pt.split(":", 1)
                run1 = p.add_run()
                run1.text = lead + ":"
                run1.font.bold = True
                run1.font.color.rgb = NAVY
                
                run2 = p.add_run()
                run2.text = detail
                run2.font.color.rgb = DARK_GRAY
            else:
                run = p.add_run()
                run.text = pt
                run.font.color.rgb = DARK_GRAY

    # ====================================================
    # Slide 1: RTG Problem Statement
    # ====================================================
    blank_layout = prs.slide_layouts[6]
    slide1 = prs.slides.add_slide(blank_layout)
    apply_slide_template(slide1, "Modernizing Port Operations: The RTG Safety Challenge", "Addressing Operational Blindspots and Safety compliance in APSEZ Terminals")
    
    # Left Column Container (Strategic Gaps)
    left_box = slide1.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(6.0), Inches(4.3)
    )
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = LIGHT_GRAY
    left_box.line.color.rgb = NAVY
    left_box.line.width = Pt(1)
    
    left_tf = left_box.text_frame
    left_tf.word_wrap = True
    left_tf.margin_left = Inches(0.2)
    left_tf.margin_right = Inches(0.2)
    left_tf.margin_top = Inches(0.2)
    
    p = left_tf.paragraphs[0]
    p.text = "Strategic Operational Gaps"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = NAVY
    p.space_after = Pt(14)
    
    gaps = [
        "Lack of Real-Time Visibility: Traditional operator training logs are manual and retro-active, delaying safety feedback cycles.",
        "Objective Assessment Gap: Difficulty in objectively measuring fatigue levels, gaze deviation, and physical crane manipulation parameters.",
        "Hazard Zone Compliance: High risks associated with ground handlers entering the Rubber-Tired Gantry (RTG) footprint undetected.",
        "Operator Fatigue Challenges: High-stress environments inside RTG cabins require continuous, automated fatigue detection loops."
    ]
    add_bullet_points(left_tf, gaps, font_size=12, space_after=8)
    
    # Right Column Container (Table)
    table_shape = slide1.shapes.add_table(5, 3, Inches(6.8), Inches(1.5), Inches(6.0), Inches(4.3))
    table = table_shape.table
    
    # Table headers
    headers = ["Operational Metric", "Legacy Manual Process", "AI Video Analytics"]
    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
            
    # Table rows
    rows_data = [
        ["Incident Response", "Post-incident review (1-2h)", "Real-time alert (<2 seconds)"],
        ["Fatigue Monitoring", "Supervisor physical checks", "Continuous AI eye-tracking"],
        ["Data Ingestion", "Manual paper logs", "Live database streaming"],
        ["Assessment Quality", "Subjective trainer rating", "Objective, telemetry scores"]
    ]
    for row_idx, row_val in enumerate(rows_data):
        for col_idx, val in enumerate(row_val):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY if row_idx % 2 == 0 else WHITE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = DARK_GRAY
                p.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.CENTER

    # Bottom Callout Box
    callout = slide1.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(12.333), Inches(0.9))
    c_tf = callout.text_frame
    c_tf.word_wrap = True
    p = c_tf.paragraphs[0]
    p.text = "Safety First Mandate: Transitioning from reactive logging to proactive, AI-driven accident prevention protects our logistics workforce and maintains APSEZ's zero-harm environment."
    p.font.name = "Segoe UI"
    p.font.size = Pt(13)
    p.font.italic = True
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER

    # ====================================================
    # Slide 2: RTG Video vs. POC
    # ====================================================
    slide2 = prs.slides.add_slide(blank_layout)
    apply_slide_template(slide2, "Vision System Architecture: RTG Video vs. Pilot POC", "Translating Raw Video Streams into Verified Safety Model Detections")
    
    # Left Column (RTG Video)
    rtg_box = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.4)
    )
    rtg_box.fill.solid()
    rtg_box.fill.fore_color.rgb = LIGHT_GRAY
    rtg_box.line.color.rgb = NAVY
    
    rtg_tf = rtg_box.text_frame
    rtg_tf.word_wrap = True
    rtg_tf.margin_left = Inches(0.3)
    rtg_tf.margin_top = Inches(0.3)
    
    p = rtg_tf.paragraphs[0]
    p.text = "Production RTG Video Analytics"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = NAVY
    p.space_after = Pt(14)
    
    rtg_points = [
        "Edge-AI Camera Feed: Multi-angle high-resolution cameras mounted on the gantry, cabin, and trolley to track operational environment.",
        "Geofencing & Person Detection: Dynamic red boundary zones mapped around the crane footprint to monitor ground-handler incursions.",
        "Operator Behavioral Monitoring: In-cabin infrared camera tracking head alignment, phone use, and micro-sleep events.",
        "Sway & Alignment Analysis: Intelligent computer vision models assessing container sway angles and chassis alignment speeds."
    ]
    add_bullet_points(rtg_tf, rtg_points, font_size=12, space_after=12)
    
    # Right Column (POC)
    poc_box = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.4)
    )
    poc_box.fill.solid()
    poc_box.fill.fore_color.rgb = LIGHT_GRAY
    poc_box.line.color.rgb = ORANGE
    
    poc_tf = poc_box.text_frame
    poc_tf.word_wrap = True
    poc_tf.margin_left = Inches(0.3)
    poc_tf.margin_top = Inches(0.3)
    
    p = poc_tf.paragraphs[0]
    p.text = "Proof of Concept (POC) Validation"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = ORANGE
    p.space_after = Pt(14)
    
    poc_points = [
        "Deployment Baseline: Tested over a 30-day window on 3 active RTG cranes in live port environments.",
        "Model Accuracy Ratings: Achieved 96.5% precision in geofence breach detections and 94.2% sensitivity for cabin fatigue alerts.",
        "Latency Performance: Average end-to-end latency of 150ms from camera frame capture to cabin alarm trigger.",
        "Industrial Grade Hardware: Verified durability of edge gateway processing units against port dust, salt spray, and vibration."
    ]
    add_bullet_points(poc_tf, poc_points, font_size=12, space_after=12)

    # ====================================================
    # Slide 3: Screenshots Placeholder
    # ====================================================
    slide3 = prs.slides.add_slide(blank_layout)
    apply_slide_template(slide3, "RTG Video Analytics & POC: Operational Interface", "Dashboard Snapshots and Vision Model Detection Output")
    
    # Large left image placeholder
    img_left = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(7.5), Inches(5.4))
    img_left.fill.solid()
    img_left.fill.fore_color.rgb = LIGHT_GRAY
    img_left.line.color.rgb = NAVY
    tf_l = img_left.text_frame
    tf_l.word_wrap = True
    tf_l.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf_l.paragraphs[0]
    p.text = "[ PLACEHOLDER: RTG Geofencing & Person Detection Visual ]\n\nInsert a screenshot showing the RTG crane base with computer vision bounding boxes around ground personnel and active safety boundary overlays."
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER
    
    # Top right smaller image placeholder
    img_tr = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.3), Inches(1.5), Inches(4.5), Inches(2.5))
    img_tr.fill.solid()
    img_tr.fill.fore_color.rgb = LIGHT_GRAY
    img_tr.line.color.rgb = NAVY
    tf_tr = img_tr.text_frame
    tf_tr.word_wrap = True
    tf_tr.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf_tr.paragraphs[0]
    p.text = "[ PLACEHOLDER: Cabin Operator Fatigue Feed ]\n\nInsert image showing eye-tracking and facial vector mappings."
    p.font.size = Pt(12)
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER
    
    # Bottom right smaller image placeholder
    img_br = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.3), Inches(4.4), Inches(4.5), Inches(2.5))
    img_br.fill.solid()
    img_br.fill.fore_color.rgb = LIGHT_GRAY
    img_br.line.color.rgb = NAVY
    tf_br = img_br.text_frame
    tf_br.word_wrap = True
    tf_br.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf_br.paragraphs[0]
    p.text = "[ PLACEHOLDER: Model Performance Graphs ]\n\nInsert precision curves, confusion matrix, or latency logs from POC."
    p.font.size = Pt(12)
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

    # ====================================================
    # Slide 4: Excel Dashboard vs. Trainer Dashboard
    # ====================================================
    slide4 = prs.slides.add_slide(blank_layout)
    apply_slide_template(slide4, "Analytics Transformation: Excel vs. Power BI Trainer Portal", "Transitioning from Legacy Manual Logging to Interactive Operational Analytics")
    
    # Left Column (Excel)
    excel_box = slide4.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(6.0), Inches(4.5)
    )
    excel_box.fill.solid()
    excel_box.fill.fore_color.rgb = LIGHT_GRAY
    excel_box.line.color.rgb = RGBColor(200, 50, 50) # Red accent for bottleneck
    
    ex_tf = excel_box.text_frame
    ex_tf.word_wrap = True
    ex_tf.margin_left = Inches(0.3)
    ex_tf.margin_top = Inches(0.3)
    
    p = ex_tf.paragraphs[0]
    p.text = "Legacy Excel Dashboards (Operational Bottleneck)"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(180, 40, 40)
    p.space_after = Pt(14)
    
    ex_pts = [
        "Manual Data Collection: Trainers spent up to 8 hours weekly compiling Excel reports from various systems.",
        "Static Visibility: No interactive filters or real-time views; charts could not display individual trainee scorecards dynamically.",
        "Siloed Systems: Attendance logs, classroom grades, and simulator telemetry existed in separate spreadsheets.",
        "Delayed Action: Reports generated post-cohort completion, preventing early interventions for struggling operators."
    ]
    add_bullet_points(ex_tf, ex_pts, font_size=12, space_after=10)
    
    # Right Column (Trainer Dashboard)
    train_box = slide4.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(6.0), Inches(4.5)
    )
    train_box.fill.solid()
    train_box.fill.fore_color.rgb = LIGHT_GRAY
    train_box.line.color.rgb = NAVY
    
    tr_tf = train_box.text_frame
    tr_tf.word_wrap = True
    tr_tf.margin_left = Inches(0.3)
    tr_tf.margin_top = Inches(0.3)
    
    p = tr_tf.paragraphs[0]
    p.text = "Modernized Power BI Trainer Dashboard"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = NAVY
    p.space_after = Pt(14)
    
    tr_pts = [
        "Automated Data Consolidation: Merges metrics from theory exams, simulator files, and attendance databases.",
        "Dynamic Interface: Trainers can toggle views between cohort-wide metrics and specific student progress tables.",
        "DAX Telemetry Measures: Calculated columns for rolling averages, target pass margins, and safety violation flags.",
        "Real-Time Diagnostic Views: Identifies learning gaps during active training cycles, improving classroom adjustment speeds."
    ]
    add_bullet_points(tr_tf, tr_pts, font_size=12, space_after=10)
    
    # Bottom savings metric banner
    metrics_banner = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6.2), Inches(12.33), Inches(0.8))
    metrics_banner.fill.solid()
    metrics_banner.fill.fore_color.rgb = NAVY
    metrics_banner.line.fill.background()
    mb_tf = metrics_banner.text_frame
    mb_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = mb_tf.paragraphs[0]
    p.text = "IMPACT: reporting overhead reduced from 8 hours/week to under 5 minutes  |  85% faster trainer decision cycles"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER

    # ====================================================
    # Slide 5: Trainee Dashboard vs. Training Overview
    # ====================================================
    slide5 = prs.slides.add_slide(blank_layout)
    apply_slide_template(slide5, "Power BI Suite: Trainee Scorecards & Cohort Overview", "Granular Diagnostics Mapped Against Strategic Training Performance Indicators")
    
    # Left Column (Trainee Scorecards)
    trainee_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.4))
    trainee_box.fill.solid()
    trainee_box.fill.fore_color.rgb = LIGHT_GRAY
    trainee_box.line.color.rgb = NAVY
    t_tf = trainee_box.text_frame
    t_tf.word_wrap = True
    t_tf.margin_left = Inches(0.3)
    t_tf.margin_top = Inches(0.3)
    
    p = t_tf.paragraphs[0]
    p.text = "Trainee Scorecard Profile"
    p.font.bold = True
    p.font.size = Pt(17)
    p.font.color.rgb = NAVY
    p.space_after = Pt(12)
    
    t_pts = [
        "Individual Growth Timeline: Tracks progression scores from the first simulator orientation run to the final certification exam.",
        "Multi-Dimensional Skill Matrix: Combines cognitive test results (Mettl) with physical simulator handling data.",
        "Automated Support Triggers: Flags trainees whose crane maneuver cycle times or safety errors exceed standard deviation limits.",
        "DAX Ranking Functionality: Automatically scores and ranks students within their cohort to motivate skill acquisition."
    ]
    add_bullet_points(t_tf, t_pts, font_size=12, space_after=10)
    
    # Highlight DAX Formula box inside the left column space
    dax_sub_box = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5.3), Inches(5.4), Inches(1.3))
    dax_sub_box.fill.solid()
    dax_sub_box.fill.fore_color.rgb = NAVY
    dax_tf = dax_sub_box.text_frame
    dax_tf.word_wrap = True
    p = dax_tf.paragraphs[0]
    p.text = "Advanced DAX Logic Utilized:"
    p.font.bold = True
    p.font.size = Pt(10)
    p.font.color.rgb = ORANGE
    p.space_after = Pt(2)
    p2 = dax_tf.add_paragraph()
    p2.text = "Trainee Rank = RANKX(ALL(Trainees), [Weighted Avg Score])\nSafety Violations Index = CALCULATE(COUNT(Telemetry[Error]), Telemetry[Severity] = \"Critical\")"
    p2.font.name = "Consolas"
    p2.font.size = Pt(8.5)
    p2.font.color.rgb = WHITE

    # Right Column (Training Overview)
    prog_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.4))
    prog_box.fill.solid()
    prog_box.fill.fore_color.rgb = LIGHT_GRAY
    prog_box.line.color.rgb = NAVY
    p_tf = prog_box.text_frame
    p_tf.word_wrap = True
    p_tf.margin_left = Inches(0.3)
    p_tf.margin_top = Inches(0.3)
    
    p = p_tf.paragraphs[0]
    p.text = "Cohort Training Analytics Overview"
    p.font.bold = True
    p.font.size = Pt(17)
    p.font.color.rgb = NAVY
    p.space_after = Pt(12)
    
    p_pts = [
        "Aggregate Performance Curves: Tracks overall average scores and simulator maneuver durations week-over-week.",
        "Curriculum Friction Analysis: Identifies specific simulator exercises (such as strong-wind hoisting) with high failure rates.",
        "Operational Readiness Index: Executive KPIs displaying the percentage of current trainees ready for deployment.",
        "Instructor Performance Benchmarks: Standardizes curriculum execution by comparing success rates across training cohorts."
    ]
    add_bullet_points(p_tf, p_pts, font_size=12, space_after=10)

    # ====================================================
    # Slide 6: Classroom Training
    # ====================================================
    slide6 = prs.slides.add_slide(blank_layout)
    apply_slide_template(slide6, "Theoretical Foundation: Structured Classroom Onboarding", "Combining Industry Safety Standards with Practical Mechanical Principles")
    
    # Left Column (Description)
    desc_box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.4))
    desc_box.fill.solid()
    desc_box.fill.fore_color.rgb = LIGHT_GRAY
    desc_box.line.color.rgb = NAVY
    d_tf = desc_box.text_frame
    d_tf.word_wrap = True
    d_tf.margin_left = Inches(0.3)
    d_tf.margin_top = Inches(0.3)
    
    p = d_tf.paragraphs[0]
    p.text = "Academic & Technical Grounding"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = NAVY
    p.space_after = Pt(14)
    
    class_pts = [
        "Core Safety Syllabus: Instructing candidates on port speed restrictions, mechanical load boundaries, and emergency protocols.",
        "Interactive Methodologies: Classroom reviews shifted from passive lectures to active case studies of real port operations.",
        "Continuous Feedback Cycles: Daily quiz results used to track comprehension before trainees advance to physical simulator runs.",
        "Unified Training Environment: Merging theory classroom materials with direct simulator cockpit exercises for unified learning."
    ]
    add_bullet_points(d_tf, class_pts, font_size=12, space_after=12)
    
    # Right Column (Screenshots placeholder)
    sc_box = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.4))
    sc_box.fill.solid()
    sc_box.fill.fore_color.rgb = LIGHT_GRAY
    sc_box.line.color.rgb = NAVY
    sc_tf = sc_box.text_frame
    sc_tf.word_wrap = True
    sc_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = sc_tf.paragraphs[0]
    p.text = "[ PLACEHOLDER: Classroom & Training Visuals ]\n\nInsert images of classroom lecture slides, training simulator setup rooms, and cohort photos."
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

    # ====================================================
    # Slide 7: Audio Analytics & LLM
    # ====================================================
    slide7 = prs.slides.add_slide(blank_layout)
    apply_slide_template(slide7, "Communication Assessment: Audio Analytics & LLM Processing", "Assessing Voice Clarity, Protocol Adherence, and Signal Discipline")
    
    # Left Column (Audio)
    aud_box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.4))
    aud_box.fill.solid()
    aud_box.fill.fore_color.rgb = LIGHT_GRAY
    aud_box.line.color.rgb = NAVY
    aud_tf = aud_box.text_frame
    aud_tf.word_wrap = True
    aud_tf.margin_left = Inches(0.3)
    aud_tf.margin_top = Inches(0.3)
    
    p = aud_tf.paragraphs[0]
    p.text = "Audio Analytics Pipeline"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = NAVY
    p.space_after = Pt(14)
    
    aud_pts = [
        "Radio Transmission Capture: Records candidate voice commands during simulated harbor container movements.",
        "Speech-to-Text Conversion: Processes raw radio audio using noise filters to generate transcriptions.",
        "Acoustic Profiling: Measures key vocal telemetry variables including words-per-minute, speech pauses, and vocal volume spikes."
    ]
    add_bullet_points(aud_tf, aud_pts, font_size=12, space_after=14)
    
    # Right Column (LLM)
    llm_box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.4))
    llm_box.fill.solid()
    llm_box.fill.fore_color.rgb = LIGHT_GRAY
    llm_box.line.color.rgb = ORANGE
    llm_tf = llm_box.text_frame
    llm_tf.word_wrap = True
    llm_tf.margin_left = Inches(0.3)
    llm_tf.margin_top = Inches(0.3)
    
    p = llm_tf.paragraphs[0]
    p.text = "LLM Evaluation Model"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = ORANGE
    p.space_after = Pt(14)
    
    llm_pts = [
        "Protocol Verification: Matches transcribed text against standard marine safety communication codes.",
        "Automated Grading Model: Generates a clear scorecard evaluating communication confidence, delay, and compliance.",
        "Actionable Feedback Generation: Summarizes communication errors (e.g., 'Failed to confirm container lock status via radio')."
    ]
    add_bullet_points(llm_tf, llm_pts, font_size=12, space_after=14)

    # ====================================================
    # Slide 8: Power Automate Workflows
    # ====================================================
    slide8 = prs.slides.add_slide(blank_layout)
    apply_slide_template(slide8, "Process Automation: Power Automate Implementation", "Streamlining Onboarding Records and Ingesting Cohort Database Files")
    
    # Left Column (Photo Renaming)
    photo_box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.4))
    photo_box.fill.solid()
    photo_box.fill.fore_color.rgb = LIGHT_GRAY
    photo_box.line.color.rgb = NAVY
    ph_tf = photo_box.text_frame
    ph_tf.word_wrap = True
    ph_tf.margin_left = Inches(0.3)
    ph_tf.margin_top = Inches(0.3)
    
    p = ph_tf.paragraphs[0]
    p.text = "Candidate Photo Renaming Workflow"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = NAVY
    p.space_after = Pt(14)
    
    ph_pts = [
        "The Bottleneck: Onboarding photos submitted with default phone filenames required manual identification and renaming.",
        "Automatic Trigger: A new folder upload to OneDrive triggers the automated flow.",
        "Database Lookup: The flow queries the central database using the candidate's name or email details.",
        "Automatic File Rename: Retrieves the candidate's unique Registration ID (REG_ID) and updates the photo filename to '[REG_ID].jpg'."
    ]
    add_bullet_points(ph_tf, ph_pts, font_size=12, space_after=12)
    
    # Right Column (Digii Dump)
    digii_box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.4))
    digii_box.fill.solid()
    digii_box.fill.fore_color.rgb = LIGHT_GRAY
    digii_box.line.color.rgb = ORANGE
    dg_tf = digii_box.text_frame
    dg_tf.word_wrap = True
    dg_tf.margin_left = Inches(0.3)
    dg_tf.margin_top = Inches(0.3)
    
    p = dg_tf.paragraphs[0]
    p.text = "Digii Data Ingestion Pipeline"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = ORANGE
    p.space_after = Pt(14)
    
    dg_pts = [
        "The Bottleneck: Manual collation and parsing of Digii portal performance exports.",
        "Automated Flow Processing: Triggered weekly to extract multi-tab Excel files from the Digii portal.",
        "Data Transformation & Clean: Automatically filter out empty rows, sanitize columns, and align date ranges.",
        "Database Append: Automatically writes cleaned data rows into the Training Database, triggering a Power BI dataset refresh."
    ]
    add_bullet_points(dg_tf, dg_pts, font_size=12, space_after=12)

    # ====================================================
    # Slide 9: Mercer Mettl Assessments
    # ====================================================
    slide9 = prs.slides.add_slide(blank_layout)
    apply_slide_template(slide9, "Pre-Screening & Suitability: Mercer Mettl Integration", "Evaluating Spatial Aptitude, Stress Tolerance, and Theoretical Competency")
    
    # Left Column (Description)
    m_box = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.4))
    m_box.fill.solid()
    m_box.fill.fore_color.rgb = LIGHT_GRAY
    m_box.line.color.rgb = NAVY
    m_tf = m_box.text_frame
    m_tf.word_wrap = True
    m_tf.margin_left = Inches(0.3)
    m_tf.margin_top = Inches(0.3)
    
    p = m_tf.paragraphs[0]
    p.text = "Objective Assessment Framework"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = NAVY
    p.space_after = Pt(14)
    
    m_pts = [
        "Cognitive Assessments: Evaluating spatial awareness, concentration, and memory metrics, which correlate to crane handling efficiency.",
        "Psychometric Testing: Profiling candidates' risk tolerance and capacity to handle stress in elevated operating cabins.",
        "Technical Evaluation Exams: Designing custom online exams testing mechanical understanding and port regulations.",
        "AI proctoring Integrity: Monitoring assessments with browser-locking and face detection to guarantee credibility."
    ]
    add_bullet_points(m_tf, m_pts, font_size=12, space_after=12)
    
    # Right Column (Screenshots placeholder)
    ms_box = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.4))
    ms_box.fill.solid()
    ms_box.fill.fore_color.rgb = LIGHT_GRAY
    ms_box.line.color.rgb = NAVY
    ms_tf = ms_box.text_frame
    ms_tf.word_wrap = True
    ms_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = ms_tf.paragraphs[0]
    p.text = "[ PLACEHOLDER: Mercer Mettl Dashboard Visuals ]\n\nInsert screenshots of proctoring reports, candidate certificate samples, and Mettl analytics screen."
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

    # ====================================================
    # Slide 10: Learnings
    # ====================================================
    slide10 = prs.slides.add_slide(blank_layout)
    apply_slide_template(slide10, "Strategic Project Review: Learnings & Business Impact", "Key Takeaways, Operational Savings, and Next Phase Roadmap")
    
    # Box 1: Technical Learnings
    b1 = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(3.8), Inches(5.4))
    b1.fill.solid()
    b1.fill.fore_color.rgb = NAVY
    b1.line.color.rgb = ORANGE
    tf1 = b1.text_frame
    tf1.word_wrap = True
    tf1.margin_left = Inches(0.2)
    tf1.margin_top = Inches(0.2)
    
    p = tf1.paragraphs[0]
    p.text = "Technical Learnings"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = ORANGE
    p.space_after = Pt(14)
    
    pts1 = [
        "Edge Device Tuning: Port environmental parameters (humidity, vibrations) require robust edge gateway configurations.",
        "ETL Automation: Direct REST API endpoints are preferred over batch file transfers to eliminate latency.",
        "DAX Code Clean: Clean and optimized measures in Power BI prevent loading issues as cohort datasets expand."
    ]
    for pt in pts1:
        p = tf1.add_paragraph()
        p.text = "• " + pt
        p.font.size = Pt(10)
        p.font.color.rgb = WHITE
        p.space_after = Pt(10)
        
    # Box 2: Business Impact (Table)
    b2 = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(1.5), Inches(4.3), Inches(5.4))
    b2.fill.solid()
    b2.fill.fore_color.rgb = LIGHT_GRAY
    b2.line.color.rgb = NAVY
    tf2 = b2.text_frame
    tf2.word_wrap = True
    tf2.margin_left = Inches(0.2)
    tf2.margin_top = Inches(0.2)
    
    p = tf2.paragraphs[0]
    p.text = "Operational Impact Metrics"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = NAVY
    p.space_after = Pt(10)
    
    # Table inside Box 2
    inner_table_shape = slide10.shapes.add_table(4, 3, Inches(4.7), Inches(2.3), Inches(3.9), Inches(4.0))
    it = inner_table_shape.table
    it_headers = ["Metric", "Before", "After"]
    for col_idx, text in enumerate(it_headers):
        cell = it.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(10)
            p.font.color.rgb = WHITE
            
    it_data = [
        ["Data Prep Time", "10 hrs/wk", "5 mins"],
        ["Candidate Selection", "Subjective", "Objective"],
        ["Incident Response", "Reactive", "Instant"]
    ]
    for r_idx, r_val in enumerate(it_data):
        for c_idx, val in enumerate(r_val):
            cell = it.cell(r_idx + 1, c_idx)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r_idx % 2 == 0 else LIGHT_GRAY
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(8.5)
                p.font.color.rgb = DARK_GRAY

    # Box 3: Future Roadmap
    b3 = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.0), Inches(1.5), Inches(3.8), Inches(5.4))
    b3.fill.solid()
    b3.fill.fore_color.rgb = NAVY
    b3.line.color.rgb = ORANGE
    tf3 = b3.text_frame
    tf3.word_wrap = True
    tf3.margin_left = Inches(0.2)
    tf3.margin_top = Inches(0.2)
    
    p = tf3.paragraphs[0]
    p.text = "Future Roadmap"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = ORANGE
    p.space_after = Pt(14)
    
    pts3 = [
        "API Integration Hub: Replaces email uploads of Excel sheets with real-time webhook endpoints.",
        "Predictive Analytics: Using trainee metrics to anticipate candidate struggle points.",
        "Audio Coach Systems: Integrating real-time tone/pace voice advice during active simulations."
    ]
    for pt in pts3:
        p = tf3.add_paragraph()
        p.text = "• " + pt
        p.font.size = Pt(10)
        p.font.color.rgb = WHITE
        p.space_after = Pt(10)

    # Save presentation
    output_path = "c:\\Users\\Admin\\Downloads\\dhaathree_website (2)\\RTG_Training_Analytics_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")
    
    # Save a copy in artifacts directory
    artifact_dir = "C:\\Users\\Admin\\.gemini\\antigravity\\brain\\4a1feceb-f5f1-4d72-bc80-6378eeeccc52"
    if os.path.exists(artifact_dir):
        artifact_path = os.path.join(artifact_dir, "RTG_Training_Analytics_Presentation.pptx")
        prs.save(artifact_path)
        print(f"Presentation copy saved to artifacts: {artifact_path}")

if __name__ == "__main__":
    create_presentation()
